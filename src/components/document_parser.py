"""Document parsing from multiple formats."""

import io
import json
from typing import List
import streamlit as st

# Import parsing libraries
import PyPDF2
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

from PIL import Image
import easyocr
import pandas as pd
from pptx import Presentation
import numpy as np

# Optional OCR improvements
try:
    import pytesseract
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False

try:
    import cv2
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

from src.models import Passage, ParseError, Document
from src.utils.logger import setup_logger, log_event
from src.utils.constants import EMBEDDING_CHUNK_SIZE, SEMANTIC_CHUNK_MIN_SIZE, SEMANTIC_CHUNK_MAX_SIZE
from src.components.pdf_ingester import PdfIngester
from src.components.vlm_captioner import VLMCaptioner

logger = setup_logger(__name__)


class DocumentParser:
    """Parse documents in multiple formats and extract text passages."""
    
    def __init__(self):
        """Initialize the parser."""
        self.ocr_reader = None  # Lazy load EasyOCR
        self.pdf_ingester = PdfIngester()
        self.vlm_captioner = None  # Lazy load (requires OPENROUTER_API_KEY)
    
    def _preprocess_image(self, image: Image.Image) -> np.ndarray:
        """Preprocess image for better OCR results.
        
        Args:
            image: PIL Image object.
        
        Returns:
            Preprocessed numpy array or original if cv2 unavailable.
        """
        if not HAS_CV2:
            return np.array(image)
        
        try:
            # Convert to numpy
            img_np = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            
            # Upscale if image is small (improves OCR)
            height, width = img_np.shape[:2]
            if width < 800:
                scale = 800 / width
                img_np = cv2.resize(img_np, None, fx=scale, fy=scale, interpolation=cv2.INTER_CUBIC)
            
            # Convert to grayscale
            gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
            
            # Apply CLAHE (Contrast Limited Adaptive Histogram Equalization)
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            enhanced = clahe.apply(gray)
            
            # Binarization with adaptive threshold
            binary = cv2.adaptiveThreshold(enhanced, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                          cv2.THRESH_BINARY, 11, 2)
            
            # Denoise
            denoised = cv2.fastNlMeansDenoising(binary, h=10)
            
            return denoised
        except Exception as e:
            logger.warning(f"Image preprocessing failed, using original: {e}")
            return np.array(image)
    
    def _ocr_image_pytesseract(self, image: Image.Image) -> str:
        """Extract text from image using Tesseract (high quality).
        
        Args:
            image: PIL Image object.
        
        Returns:
            Extracted text or None if pytesseract unavailable.
        """
        if not HAS_PYTESSERACT:
            return None
        
        try:
            # Preprocess image for better results
            preprocessed = self._preprocess_image(image)
            
            # Convert back to PIL for pytesseract
            preprocessed_pil = Image.fromarray(preprocessed)
            
            # Run Tesseract with high config
            text = pytesseract.image_to_string(
                preprocessed_pil,
                config='--psm 6 --oem 3'  # PSM 6: Assume single block of text, OEM 3: both modes
            )
            
            return text.strip()
        except Exception as e:
            logger.debug(f"Pytesseract OCR failed, falling back to EasyOCR: {e}")
            return None
    
    def parse(self, file_bytes: bytes, filename: str, file_format: str, document_id: str) -> Document:
        """Parse a document and return extracted passages.
        
        Args:
            file_bytes: Raw file bytes.
            filename: Original filename.
            file_format: Format type (pdf, image, excel, powerpoint, video).
            document_id: Unique document ID.
        
        Returns:
            Document object with parsed passages.
        
        Raises:
            ParseError: If parsing fails.
        """
        try:
            if file_format == "pdf":
                # PDFs go through the dedicated ingestion pipeline (PyMuPDF4LLM + VLM
                # diagram captions); header-aware chunking with overlap happens
                # downstream in the indexing pipeline (see Chunker), so page-level
                # passages are returned here as-is.
                passages = self._parse_pdf_pipeline(file_bytes)
            elif file_format == "image":
                passages = self._parse_image(file_bytes)
            elif file_format == "excel":
                passages = self._parse_excel(file_bytes)
            elif file_format == "powerpoint":
                passages = self._parse_powerpoint(file_bytes)
            elif file_format == "video":
                passages = self._parse_video(file_bytes)
            else:
                raise ParseError(f"Unsupported format: {file_format}")
            
            # Legacy formats are chunked here; PDFs are chunked downstream by Chunker.
            if file_format != "pdf":
                passages = self._chunk_passages(passages)
            
            # Assign document_id to all passages
            for p in passages:
                p.document_id = document_id
            
            log_event("parse_success", filename=filename, format=file_format, passage_count=len(passages))
            
            return Document(
                document_id=document_id,
                filename=filename,
                file_format=file_format,
                file_hash="",
                parsed_successfully=True,
                passages=passages,
            )
        
        except Exception as e:
            logger.exception(f"Parse failed for {filename}")
            return Document(
                document_id=document_id,
                filename=filename,
                file_format=file_format,
                file_hash="",
                parsed_successfully=False,
                passages=[],
                parse_error_message=str(e),
            )
    
    def _parse_pdf_pipeline(self, file_bytes: bytes) -> List[Passage]:
        """Parse a PDF via the diagnostic ingestion pipeline: PyMuPDF4LLM for
        header-aware markdown text (with page metadata) + VLM captions for
        embedded diagrams. Falls back to the legacy pdfplumber/PyPDF2 + OCR
        parser if the new pipeline fails (e.g. malformed PDF).

        Args:
            file_bytes: Raw PDF bytes.

        Returns:
            List of Passage objects (one per page, plus one per captioned diagram).
        """
        try:
            ingested = self.pdf_ingester.ingest(file_bytes)
            passages: List[Passage] = []

            for page in ingested["pages"]:
                text = page["markdown"].strip()
                if text:
                    passages.append(
                        Passage(
                            text=text,
                            section=f"Page {page['page']}",
                            document_id="",
                            passage_index=len(passages),
                            is_diagram=False,
                        )
                    )

            for image in ingested["images"]:
                caption = self._caption_diagram(image["bytes"])
                if caption:
                    passages.append(
                        Passage(
                            text=caption,
                            section=f"Page {image['page']} - Diagram {image['index']}",
                            document_id="",
                            passage_index=len(passages),
                            image_bytes=image["bytes"],
                            is_diagram=True,
                        )
                    )

            if not passages:
                raise ParseError("No text or diagrams extracted from PDF")

            return passages
        except Exception as e:
            logger.warning(f"PDF ingestion pipeline failed, falling back to legacy parser: {e}")
            return self._parse_pdf(file_bytes)

    def _caption_diagram(self, image_bytes: bytes) -> str:
        """Caption a diagram image via the VLM, lazily initializing the captioner."""
        try:
            if self.vlm_captioner is None:
                self.vlm_captioner = VLMCaptioner()
            return self.vlm_captioner.caption(image_bytes)
        except Exception as e:
            logger.warning(f"Diagram captioning unavailable: {e}")
            return ""

    def _parse_pdf(self, file_bytes: bytes) -> List[Passage]:
        """Parse a PDF file and extract text + images with hybrid embeddings.
        
        Strategy:
        1. Extract text from pages
        2. Extract images and store them as separate diagram passages
        3. Use improved OCR (Tesseract with preprocessing) on images
        4. Create dual embeddings: image embedding + OCR text embedding
        
        Args:
            file_bytes: Raw PDF bytes.
        
        Returns:
            List of Passage objects (text + diagram passages).
        
        Raises:
            ParseError: If PDF cannot be parsed.
        """
        passages = []
        
        try:
            # Try pdfplumber first (better extraction + image support)
            if pdfplumber:
                pdf = pdfplumber.open(io.BytesIO(file_bytes))
                for page_num, page in enumerate(pdf.pages, 1):
                    # Extract text
                    text = page.extract_text()
                    if text and text.strip():
                        passages.append(
                            Passage(
                                text=text,
                                section=f"Page {page_num}",
                                document_id="",
                                passage_index=len(passages),
                                is_diagram=False,
                            )
                        )
                    
                    # Extract and process images (hybrid approach)
                    try:
                        if hasattr(page, 'images') and page.images:
                            for img_idx, img in enumerate(page.images, 1):
                                try:
                                    # Extract image object
                                    if hasattr(img, 'stream'):
                                        img_bytes = img.stream.get_rawdata()
                                        img_obj = Image.open(io.BytesIO(img_bytes))
                                    else:
                                        cropped = page.crop((img['x0'], img['top'], img['x1'], img['bottom']))
                                        img_obj = cropped.to_image()
                                    
                                    # Store image bytes
                                    img_bytes_io = io.BytesIO()
                                    img_obj.save(img_bytes_io, format='PNG')
                                    img_bytes_final = img_bytes_io.getvalue()
                                    
                                    # OCR with improved Tesseract (high quality for technical diagrams)
                                    ocr_text = self._ocr_image_pytesseract(img_obj)
                                    
                                    # Fallback to EasyOCR if Tesseract fails
                                    if not ocr_text:
                                        if self.ocr_reader is None:
                                            logger.info("Initializing EasyOCR reader...")
                                            self.ocr_reader = easyocr.Reader(['en'], gpu=False)
                                        result = self.ocr_reader.readtext(img_obj)
                                        ocr_text = '\n'.join([text for (_, text, _) in result])
                                    
                                    # Create diagram passage with both image bytes and OCR text
                                    if ocr_text.strip() or img_bytes_final:
                                        passages.append(
                                            Passage(
                                                text=f"DIAGRAM {img_idx} (Page {page_num}):\n{ocr_text.strip()}",
                                                section=f"Page {page_num} - Diagram {img_idx}",
                                                document_id="",
                                                passage_index=len(passages),
                                                image_bytes=img_bytes_final,
                                                is_diagram=True,
                                            )
                                        )
                                        logger.info(f"Extracted diagram {img_idx} on page {page_num} (size: {len(img_bytes_final)} bytes)")
                                
                                except Exception as img_err:
                                    logger.warning(f"Failed to process image {img_idx} on page {page_num}: {img_err}")
                    
                    except Exception as img_extraction_err:
                        logger.debug(f"Image extraction failed for page {page_num}: {img_extraction_err}")
                
                return passages
        except Exception as e:
            logger.debug(f"pdfplumber failed, trying PyPDF2: {e}")
        
        # Fallback to PyPDF2 (text only, no image support)
        try:
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    passages.append(
                        Passage(
                            text=text,
                            section=f"Page {page_num}",
                            document_id="",
                            passage_index=page_num - 1,
                            is_diagram=False,
                        )
                    )
            return passages
        except Exception as e:
            raise ParseError(f"Failed to parse PDF: {str(e)}")
    
    def _parse_image(self, file_bytes: bytes) -> List[Passage]:
        """Parse an image file using OCR.
        
        Args:
            file_bytes: Raw image bytes.
        
        Returns:
            List of Passage objects with OCR text.
        
        Raises:
            ParseError: If OCR fails.
        """
        try:
            # Lazy load OCR reader
            if self.ocr_reader is None:
                logger.info("Initializing EasyOCR reader...")
                self.ocr_reader = easyocr.Reader(['en'])
            
            # Convert bytes to image
            image = Image.open(io.BytesIO(file_bytes))
            
            # Run OCR
            results = self.ocr_reader.readtext(image)
            
            # Extract text
            text = "\n".join([result[1] for result in results])
            
            if not text or not text.strip():
                raise ParseError("No text found in image (OCR returned empty result)")
            
            return [
                Passage(
                    text=text,
                    section="Image",
                    document_id="",
                    passage_index=0,
                )
            ]
        except Exception as e:
            raise ParseError(f"Failed to parse image with OCR: {str(e)}")
    
    def _parse_excel(self, file_bytes: bytes) -> List[Passage]:
        """Parse an Excel file.
        
        Args:
            file_bytes: Raw Excel bytes.
        
        Returns:
            List of Passage objects with cell contents.
        
        Raises:
            ParseError: If Excel cannot be parsed.
        """
        try:
            xls = pd.ExcelFile(io.BytesIO(file_bytes))
            passages = []
            
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(io.BytesIO(file_bytes), sheet_name=sheet_name)
                
                # Convert dataframe to text
                text = f"Sheet: {sheet_name}\n"
                text += df.to_string()
                
                if text.strip():
                    passages.append(
                        Passage(
                            text=text,
                            section=f"Sheet: {sheet_name}",
                            document_id="",
                            passage_index=len(passages),
                        )
                    )
            
            if not passages:
                raise ParseError("No content found in Excel file")
            
            return passages
        except Exception as e:
            raise ParseError(f"Failed to parse Excel file: {str(e)}")
    
    def _parse_powerpoint(self, file_bytes: bytes) -> List[Passage]:
        """Parse a PowerPoint file.
        
        Args:
            file_bytes: Raw PowerPoint bytes.
        
        Returns:
            List of Passage objects with slide contents.
        
        Raises:
            ParseError: If PowerPoint cannot be parsed.
        """
        try:
            prs = Presentation(io.BytesIO(file_bytes))
            passages = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            text_content.append(text)
                
                if text_content:
                    passage_text = "\n".join(text_content)
                    passages.append(
                        Passage(
                            text=passage_text,
                            section=f"Slide {slide_num}",
                            document_id="",
                            passage_index=slide_num - 1,
                        )
                    )
            
            if not passages:
                raise ParseError("No content found in PowerPoint file")
            
            return passages
        except Exception as e:
            raise ParseError(f"Failed to parse PowerPoint file: {str(e)}")
    
    def _parse_video(self, file_bytes: bytes) -> List[Passage]:
        """Parse a video file (extract metadata, captions if available).
        
        Note: Full speech-to-text transcription requires OpenAI Whisper API
        and is expensive. For v1, we support captions only.
        
        Args:
            file_bytes: Raw video bytes.
        
        Returns:
            List of Passage objects with video metadata/captions.
        
        Raises:
            ParseError: If video cannot be parsed.
        """
        try:
            # For v1, we'll return a simple message
            # In v2, integrate OpenAI Whisper API for transcription
            return [
                Passage(
                    text="Video file uploaded. Speech-to-text transcription is available with OpenAI Whisper API integration (v2).",
                    section="Video Metadata",
                    document_id="",
                    passage_index=0,
                )
            ]
        except Exception as e:
            raise ParseError(f"Failed to process video file: {str(e)}")
    
    def _chunk_passages(self, passages: List[Passage]) -> List[Passage]:
        """Chunk passages semantically by section headers and size limits.
        
        Strategy:
        1. Split by section headers (##, ###, etc.)
        2. Keep diagram+text together within sections
        3. Respect min/max size limits
        4. Preserve page context in section names
        
        Args:
            passages: List of passages, some may be large.
        
        Returns:
            List of passages, chunked semantically.
        """
        import re
        
        chunked = []
        for passage in passages:
            # Pattern for headers: lines starting with # or ALL CAPS followed by colon
            header_pattern = r'^(#{1,3}\s+.+|[A-Z\s]+:)$'
            lines = passage.text.split('\n')
            
            sections = []
            current_section = None
            current_text = []
            
            for line in lines:
                # Check if line is a section header
                if re.match(header_pattern, line.strip()):
                    # Save previous section if exists
                    if current_section and current_text:
                        section_text = '\n'.join(current_text).strip()
                        if len(section_text) > SEMANTIC_CHUNK_MIN_SIZE:
                            sections.append((current_section, section_text))
                        else:
                            # Merge small section with next one
                            current_text = [section_text]
                    
                    # Start new section
                    current_section = line.strip()
                    current_text = [current_section]
                else:
                    # Add line to current section
                    if current_text or line.strip():  # Skip leading empty lines
                        current_text.append(line)
            
            # Add final section
            if current_section and current_text:
                section_text = '\n'.join(current_text).strip()
                if len(section_text) > SEMANTIC_CHUNK_MIN_SIZE:
                    sections.append((current_section, section_text))
            
            # If no sections found (no headers), chunk by size
            if not sections:
                sections = [("Page", passage.text)]
            
            # Now create passages from semantic sections, respecting max size
            for idx, (section_name, section_text) in enumerate(sections):
                # Further split if section exceeds max size
                if len(section_text) > SEMANTIC_CHUNK_MAX_SIZE:
                    # Split by sentences or paragraphs
                    sub_chunks = self._split_section(section_text, SEMANTIC_CHUNK_MAX_SIZE)
                    for sub_idx, sub_chunk in enumerate(sub_chunks):
                        if sub_chunk.strip():
                            chunked.append(
                                Passage(
                                    text=sub_chunk,
                                    section=f"{passage.section} > {section_name} ({sub_idx + 1}/{len(sub_chunks)})",
                                    document_id=passage.document_id,
                                    passage_index=len(chunked),
                                )
                            )
                else:
                    # Use section as-is
                    if section_text.strip():
                        chunked.append(
                            Passage(
                                text=section_text,
                                section=f"{passage.section} > {section_name}",
                                document_id=passage.document_id,
                                passage_index=len(chunked),
                            )
                        )
            
            logger.info(f"Semantic chunking: '{passage.section}' → {len([s for s in sections if len(s[1]) > SEMANTIC_CHUNK_MIN_SIZE])} semantic chunks")
        
        return chunked if chunked else passages
    
    def _split_section(self, text: str, max_size: int) -> List[str]:
        """Split a large section by paragraph/sentence boundaries.
        
        Args:
            text: Section text to split
            max_size: Maximum size per chunk
        
        Returns:
            List of chunks, each ≤ max_size
        """
        # Split by double newlines (paragraphs)
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = []
        current_size = 0
        
        for para in paragraphs:
            para_size = len(para)
            if current_size + para_size > max_size and current_chunk:
                # Flush current chunk
                chunks.append('\n\n'.join(current_chunk))
                current_chunk = [para]
                current_size = para_size
            else:
                current_chunk.append(para)
                current_size += para_size + 2  # +2 for \n\n
        
        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))
        
        return chunks
